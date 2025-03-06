#pip install langchain python-pptx pandas matplotlib seaborn requests openai python-docx

from typing import List, Optional
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from langchain.agents import initialize_agent, AgentType
from langchain_core.tools import tool
from langchain.document_loaders import CSVLoader
from langchain_community.document_loaders.excel import UnstructuredExcelLoader
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders import WebBaseLoader
from langchain_openai import ChatOpenAI  # 导入 ChatOpenAI 类，用于与 OpenAI API 交互
from langchain_core.tools import tool
from dotenv import load_dotenv  # 导入 dotenv 库，用于加载环境变量

import seaborn as sns
import os  # 导入 os 库，用于操作系统相关功能

# 加载 .env 文件中的环境变量
load_dotenv()

# 获取 OpenAI API 的 base URL
base_url = "https://api.deepseek.com"
api_key = os.getenv("DEEPSEEK_API_KEY")
model_name = "deepseek-chat"

llm = ChatOpenAI(base_url=base_url,
                 model=model_name,
                 api_key=api_key)  # 创建 ChatOpenAI 实例




class ChartGenerator:
    @staticmethod
    def create_chart(data: pd.DataFrame, chart_type: str, x_column: str, y_column: str, 
                     title: str, output_path: str) -> str:
        """
        生成各种类型的图表
        """
        plt.figure(figsize=(10, 6))
        if chart_type == "bar":
            sns.barplot(data=data, x=x_column, y=y_column)
        elif chart_type == "line":
            sns.lineplot(data=data, x=x_column, y=y_column)
        elif chart_type == "pie":
            plt.pie(data[y_column], labels=data[x_column], autopct='%1.1f%%')
        elif chart_type == "scatter":
            sns.scatterplot(data=data, x=x_column, y=y_column)
        plt.rcParams['font.sans-serif']=['SimHei']   # 用黑体显示中文
        plt.rcParams['axes.unicode_minus']=False     # 正常显示负号
        plt.title(title)
        plt.tight_layout()
        plt.savefig(output_path)
        plt.close()
        return output_path

@tool
def generate_chart(data_path: str, chart_type: str, x_column: str, y_column: str, 
                   title: str, output_path: str = "chart.png") -> str:
    """
    生成图表的工具
    Args:
        data_path: 数据文件路径
        chart_type: 图表类型 ('bar', 'line', 'pie', 'scatter')
        x_column: X轴列名
        y_column: Y轴列名
        title: 图表标题
        output_path: 输出图片路径
    Returns:
        str: 生成的图表文件路径
    """
    data = pd.read_csv(data_path) if data_path.endswith('.csv') else pd.read_excel(data_path)
    try:
        result = ChartGenerator.create_chart(data, chart_type, x_column, y_column, title, output_path)
    except Exception as e:
        result = f"生成图表失败: {e}"
        
    return result

@tool
def generate_and_save_ppt(ppt_content: List[dict], output_filename: Optional[str] = "generated_ppt.pptx") -> str:
    """
    生成PPT的工具
    Args:
        ppt_content: PPT内容列表，每个字典代表一页幻灯片的内容, 包含以下字段:
            - slide_type: 幻灯片类型 ('title', 'text_only', 'text_and_image')
            - title: 幻灯片标题
            - subtitle: 幻灯片副标题
            - content: 幻灯片内容列表
            - image_path: 图片路径
        output_filename: 输出文件名
    Returns:
        str: 生成的PPT文件路径
    """
    prs = Presentation()
    for slide_info in ppt_content:
        slide_type = slide_info["slide_type"]
        if slide_type == "title":
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_info["title"]
            subtitle.text = slide_info.get("subtitle", "")
        
        elif slide_type == "text_only":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = slide_info["title"]
            content.text = "\n".join(slide_info.get("content", []))
            
        elif slide_type == "text_and_image":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_info["title"]
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
            content_frame = content_box.text_frame
            for paragraph in slide_info.get("content", []):
                content_frame.add_paragraph().text = paragraph
                
            image_path = slide_info.get("image_path")
            if image_path:
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                image_width = Inches(4)
                image_height = Inches(3)
                image_left = slide_width - image_width - Inches(0.5)
                image_top = (slide_height - image_height) / 2
                slide.shapes.add_picture(image_path, image_left, image_top, 
                                      width=image_width, height=image_height)
    
    prs.save(output_filename)
    return output_filename

@tool
def file_loader(file_path: str) -> str:
    """
    文件加载工具
    Args:
        file_path: 文件路径
    Returns:
        str: 文件内容
    """
    file_ext = file_path.split('.')[-1].lower()
    if file_ext == 'csv':
        loader = CSVLoader(file_path)
    elif file_ext in ['xls', 'xlsx']:
        loader = UnstructuredExcelLoader(file_path)
    elif file_ext in ['doc', 'docx']:
        loader = Docx2txtLoader(file_path)
    elif file_path.startswith("http"):
        loader = WebBaseLoader(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_ext}")
    
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def create_ppt_agent(llm):
    """
    创建PPT生成Agent
    """
    tools = [generate_chart, generate_and_save_ppt, file_loader]
    
    agent = initialize_agent(
        tools,
        llm,
        agent=AgentType.STRUCTURED_CHAT_ZERO_SHOT_REACT_DESCRIPTION,
        verbose=True
    )
    
    return agent

# 使用示例
if __name__ == "__main__":
    agent = create_ppt_agent(llm)
    
    # 示例调用
    result = agent.run({
        "input": """
        使用 ./sales_data1.csv 文件创建一个销售分析PPT，包括以下内容：
        1. 按月份绘制销售额的折线图
        2. 分析销售额最高的日期
        3. 按产品类别绘制销售额的柱状图
        4. 异常值检测
        5. 趋势分析
        6. 商业建议
        7. 数据的关键特征
        8. 主要统计指标
        9. 潜在的数据模式
        等等
        内容请使用金字塔原理组织。
        """
    })
    
    print(f"PPT已生成: {result}")

